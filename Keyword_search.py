import re
from pathlib import Path
from collections import Counter, defaultdict

import numpy as np
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

from sklearn.feature_extraction.text import TfidfVectorizer, ENGLISH_STOP_WORDS
from sklearn.metrics.pairwise import cosine_similarity


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".xlsx", ".pptx"}


# -----------------------------
# Text Extraction
# -----------------------------
def extract_text_from_pdf(file_path: Path) -> str:
    text_parts = []
    reader = PdfReader(str(file_path))
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n".join(text_parts)


def extract_text_from_docx(file_path: Path) -> str:
    doc = Document(str(file_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_txt(file_path: Path) -> str:
    try:
        return file_path.read_text(encoding="utf-8", errors="strict")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="latin-1", errors="replace")


def extract_text_from_xlsx(file_path: Path) -> str:
    """
    Extract text from all sheets and cells in Excel (.xlsx).
    """
    parts = []
    xls = pd.ExcelFile(str(file_path), engine="openpyxl")
    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet, engine="openpyxl", dtype=str).fillna("")
        sheet_text = "\n".join(
            " ".join(str(v).strip() for v in row if str(v).strip())
            for row in df.values.tolist()
        ).strip()
        if sheet_text:
            parts.append(f"[SHEET: {sheet}]\n{sheet_text}")
    return "\n\n".join(parts)


def extract_text_from_pptx(file_path: Path) -> str:
    """
    Extract text from slides/shapes in PowerPoint (.pptx).
    """
    prs = Presentation(str(file_path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_parts.append(t)
        if slide_parts:
            parts.append(f"[SLIDE {i}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# -----------------------------
# Load ONLY from given folder (non-recursive)
# -----------------------------
def load_documents(folder_path: str):
    """
    Loads supported documents ONLY from the given folder path.
    Does NOT include subfolders.
    """
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        raise ValueError("❌ The path you entered is not a valid folder. Please try again.")

    documents = []
    for file_path in folder.glob("*"):  # only this folder
        if not file_path.is_file():
            continue

        ext = file_path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue

        try:
            if ext == ".pdf":
                text = extract_text_from_pdf(file_path)
            elif ext == ".docx":
                text = extract_text_from_docx(file_path)
            elif ext == ".txt":
                text = extract_text_from_txt(file_path)
            elif ext == ".xlsx":
                text = extract_text_from_xlsx(file_path)
            elif ext == ".pptx":
                text = extract_text_from_pptx(file_path)
            else:
                continue

            text = normalize_text(text)
            if len(text.split()) >= 15:
                documents.append({
                    "path": str(file_path),
                    "name": file_path.name,
                    "text": text
                })

        except Exception as e:
            print(f"⚠️ Skipping file (error reading): {file_path.name} -> {e}")

    return documents


# -----------------------------
# Chunking + Index
# -----------------------------
def chunk_text(text: str, chunk_size=180, overlap=40):
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        if end == len(words):
            break
        start = end - overlap
    return chunks


def build_search_index(documents):
    chunk_records, chunk_texts = [], []
    for doc in documents:
        for i, ch in enumerate(chunk_text(doc["text"])):
            chunk_records.append({
                "path": doc["path"],
                "name": doc["name"],
                "chunk_id": i,
                "text": ch
            })
            chunk_texts.append(ch)

    if not chunk_texts:
        raise ValueError("❌ No readable text found in the supported documents.")

    vectorizer = TfidfVectorizer(
        lowercase=True,
        stop_words="english",
        ngram_range=(1, 2),
        max_features=50000
    )
    chunk_matrix = vectorizer.fit_transform(chunk_texts)
    return chunk_records, vectorizer, chunk_matrix


# -----------------------------
# Multi-keyword Search
# -----------------------------
def parse_keywords(user_input: str):
    """
    Accepts: "a,b,c" OR "a b c"
    Returns: list of terms
    """
    s = user_input.strip()
    if not s:
        return []
    if "," in s or ";" in s:
        parts = re.split(r"[,;]+", s)
        return [p.strip() for p in parts if p.strip()]
    return [t.strip() for t in s.split() if t.strip()]


def keyword_filter(text: str, keywords, mode="OR"):
    t = text.lower()
    hits = [k.lower() in t for k in keywords]
    return all(hits) if mode == "AND" else any(hits)


def build_query_from_keywords(keywords, mode="OR"):
    """
    AND mode repeats keywords to increase TF-IDF emphasis.
    """
    if mode == "AND":
        return " ".join(k for k in keywords for _ in range(2))
    return " ".join(keywords)


def search_relevant_chunks(query, keywords, mode, chunk_records, vectorizer, chunk_matrix,
                           top_k=10, min_score=0.03):

    query_vec = vectorizer.transform([query])
    scores = cosine_similarity(chunk_matrix, query_vec).ravel()

    ranked = np.argsort(scores)[::-1]
    results = []

    for idx in ranked:
        score = float(scores[idx])
        if score < min_score:
            break

        rec = dict(chunk_records[idx])
        if not keyword_filter(rec["text"], keywords, mode=mode):
            continue

        rec["score"] = score
        results.append(rec)
        if len(results) >= top_k:
            break

    return results


# -----------------------------
# Query-based Stylized EXTRACTIVE Summary (File-only)
# -----------------------------
def split_sentences(text: str):
    text = text.replace("\n", " ").strip()
    sentences = re.split(r'(?<=[.!?])\s+', text)
    # keep meaningful sentences
    return [s.strip() for s in sentences if len(s.split()) >= 6]


def tokenize(text: str):
    return re.findall(r"\b[a-zA-Z][a-zA-Z0-9_-]*\b", text.lower())


def normalize_scores(values):
    values = np.array(values, dtype=float)
    if len(values) == 0:
        return values
    mn, mx = values.min(), values.max()
    if mx == mn:
        return np.zeros_like(values) if mx == 0 else np.ones_like(values)
    return (values - mn) / (mx - mn)


def sentence_records_from_results(results):
    """
    Build sentence-level records while keeping source info.
    Ensures summaries can be cited and remain file-only.
    """
    sent_recs = []
    for r in results:
        sents = split_sentences(r["text"])
        for j, s in enumerate(sents):
            sent_recs.append({
                "sentence": s,
                "path": r["path"],
                "name": r["name"],
                "chunk_id": r["chunk_id"],
                "sent_pos": j
            })
    return sent_recs


def rank_sentences(sent_recs, query):
    """
    Rank sentences by query similarity + frequency + query term coverage.
    Output remains strictly extractive (sentences are verbatim).
    """
    if not sent_recs:
        return []

    sentences = [sr["sentence"] for sr in sent_recs]

    tfidf = TfidfVectorizer(lowercase=True, stop_words="english", ngram_range=(1, 2))
    X = tfidf.fit_transform(sentences + [query])

    sim = cosine_similarity(X[:-1], X[-1]).ravel()

    # Frequency model over ALL candidate sentences (still extractive)
    all_text = " ".join(sentences)
    words = [w for w in tokenize(all_text) if w not in ENGLISH_STOP_WORDS and len(w) > 2]
    freq = Counter(words)
    max_freq = max(freq.values()) if freq else 1

    query_terms = set(tokenize(query)) - set(ENGLISH_STOP_WORDS)

    freq_scores, cov_scores = [], []
    for s in sentences:
        sw = tokenize(s)
        f = sum(freq.get(w, 0) / max_freq for w in sw) / max(len(sw), 1) if sw else 0
        c = len(set(sw) & query_terms) / max(len(query_terms), 1) if query_terms else 0
        freq_scores.append(f)
        cov_scores.append(c)

    sim = normalize_scores(sim)
    freq_scores = normalize_scores(freq_scores)
    cov_scores = normalize_scores(cov_scores)

    final = 0.65 * sim + 0.20 * freq_scores + 0.15 * cov_scores

    ranked_idx = np.argsort(final)[::-1]
    ranked = []
    for idx in ranked_idx:
        rec = dict(sent_recs[idx])
        rec["final_score"] = float(final[idx])
        ranked.append(rec)
    return ranked


def format_citation(rec):
    # Compact, traceable citation
    return f"[{Path(rec['path']).name} | chunk {rec['chunk_id']}]"


def build_stylized_extractive_summary(query, results,
                                      tldr_sentences=2,
                                      key_points=6,
                                      evidence_per_doc=2):
    """
    Internet-inspired structure, but 100% file-only:
    - TL;DR (top sentences)
    - Key Points (bullets)
    - Evidence by file (snippets)
    Every line is a verbatim sentence from the documents.
    """
    if not results:
        return None

    sent_recs = sentence_records_from_results(results)
    ranked = rank_sentences(sent_recs, query)

    if not ranked:
        return {
            "tldr": ["No meaningful extractive summary could be generated from the matched text."],
            "key_points": [],
            "evidence": [],
            "per_doc": []
        }

    # Deduplicate identical sentences (keep best scoring)
    seen = set()
    unique_ranked = []
    for r in ranked:
        key = r["sentence"].strip().lower()
        if key in seen:
            continue
        seen.add(key)
        unique_ranked.append(r)

    # TL;DR and Key Points
    tldr = unique_ranked[:tldr_sentences]
    keyp = unique_ranked[tldr_sentences:tldr_sentences + key_points]

    # Evidence grouped by document
    grouped = defaultdict(list)
    for r in unique_ranked:
        grouped[r["path"]].append(r)

    evidence = []
    per_doc = []
    for path, items in grouped.items():
        items_sorted = sorted(items, key=lambda x: x["final_score"], reverse=True)
        top_evidence = items_sorted[:evidence_per_doc]
        evidence.extend(top_evidence)

        # per-doc highlight (keep it short + extractive)
        doc_best = items_sorted[:min(4, len(items_sorted))]
        per_doc.append({
            "path": path,
            "name": Path(path).name,
            "best_score": max(i["final_score"] for i in items_sorted),
            "highlights": doc_best
        })

    per_doc.sort(key=lambda x: x["best_score"], reverse=True)

    return {
        "tldr": tldr,
        "key_points": keyp,
        "evidence": evidence,
        "per_doc": per_doc
    }


# -----------------------------
# CLI Printing (Stylized)
# -----------------------------
def print_results(keywords, mode, summary_result):
    if not summary_result:
        print("\n❌ No relevant content found for your keywords.")
        return

    print("\n" + "=" * 90)
    print(f"✅ KEYWORDS ENTERED ({mode}): {keywords}")
    print("=" * 90)

    print("\n🧾 NOTE: Every line below is a direct extracted sentence from your files (no paraphrasing).")

    # TL;DR
    print("\n🔥 TL;DR\n")
    for r in summary_result["tldr"]:
        print(f"- {r['sentence']} {format_citation(r)}")

    # Key Points
    if summary_result["key_points"]:
        print("\n📌 KEY POINTS\n")
        for r in summary_result["key_points"]:
            print(f"- {r['sentence']} {format_citation(r)}")

    # Evidence
    if summary_result["evidence"]:
        print("\n🧩 EVIDENCE SNIPPETS (Top supporting lines)\n")
        for r in summary_result["evidence"][:10]:
            print(f"- {r['sentence']} {format_citation(r)}")

    # Per-doc highlights
    print("\n" + "-" * 90)
    print("📁 PER-FILE HIGHLIGHTS (most relevant extracts)")
    print("-" * 90)

    for i, d in enumerate(summary_result["per_doc"], start=1):
        print(f"\n{i}. File: {d['path']}")
        print(f"   Relevance Score (sentence-based): {d['best_score']:.4f}")
        for h in d["highlights"]:
            print(f"   - {h['sentence']} {format_citation(h)}")


# -----------------------------
# Main
# -----------------------------
def main():
    print("\n" + "=" * 90)
    print("DOCUMENT SEARCH + SUMMARY TOOL (PDF / DOCX / TXT / XLSX / PPTX)")
    print("=" * 90)
    print("This tool will scan ONLY the folder you provide (NO subfolders).")
    print("Then it will find relevant text using your keywords and show stylized extractive summaries.\n")

    # 1) Folder Input
    while True:
        folder_path = input("👉 Enter the folder path where your documents are located: ").strip().strip('"')
        try:
            documents = load_documents(folder_path)
            break
        except ValueError as ve:
            print(str(ve))

    if not documents:
        print("\n❌ No supported documents found in that folder.")
        print("Supported formats: .pdf, .docx, .txt, .xlsx, .pptx")
        return

    print(f"\n✅ Found {len(documents)} supported documents in the folder.")
    print("⏳ Building search index... (this may take a few seconds)")

    chunk_records, vectorizer, chunk_matrix = build_search_index(documents)
    print(f"✅ Index built successfully. Total chunks indexed: {len(chunk_records)}\n")

    # 2) Keyword Input Loop
    while True:
        print("\n" + "-" * 90)
        raw = input("🔎 Enter keyword(s) to search (example: invoice, delay, penalty) OR type 'exit': ").strip()
        if raw.lower() == "exit":
            print("\n✅ Exiting. Thank you!")
            break

        keywords = parse_keywords(raw)
        if not keywords:
            print("❌ Please enter at least one keyword.")
            continue

        # 3) Mode Input
        mode_input = input("⚙️ Choose search mode: AND (all keywords) / OR (any keyword). Default is OR: ").strip().upper()
        mode = mode_input if mode_input in {"AND", "OR"} else "OR"

        query = build_query_from_keywords(keywords, mode=mode)

        # Search
        results = search_relevant_chunks(
            query=query,
            keywords=keywords,
            mode=mode,
            chunk_records=chunk_records,
            vectorizer=vectorizer,
            chunk_matrix=chunk_matrix,
            top_k=10,
            min_score=0.03
        )

        # Stylized extractive summary (file-only)
        summary_result = build_stylized_extractive_summary(query, results)
        print_results(keywords, mode, summary_result)


if __name__ == "__main__":
    main()
